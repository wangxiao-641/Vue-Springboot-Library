from pathlib import Path
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = Path(__file__).resolve().parent / "图书馆需求变化前后对比答辩版.pptx"
W, H = 13.333, 7.5
TOTAL = 16

NAVY = "102A43"
NAVY_2 = "0B172A"
INK = "263746"
MUTED = "6B7C88"
CREAM = "F7F1E8"
WHITE = "FFFFFF"
LINE = "D9E3E1"
TEAL = "0F766E"
TEAL_LIGHT = "E6FFFA"
BRICK = "B44D3C"
BRICK_LIGHT = "F8E8DF"
PURPLE = "8B7CFF"
PURPLE_LIGHT = "24204F"
CYAN = "28D7F2"
GREEN = "47E6A8"
ORANGE = "C87845"
RED = "C63D3D"

FONT = "Microsoft YaHei"
prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

THEME_IMAGES = {
    "login": {
        "atlas": ROOT / "acceptance/screenshots/theme-v2/atlas-login.png",
        "academy": ROOT / "acceptance/screenshots/theme-v2/academy-login.png",
        "command": ROOT / "acceptance/screenshots/theme-v2/command-login.png",
    },
    "dashboard": {
        "atlas": ROOT / "acceptance/screenshots/theme-v2/atlas-dashboard.png",
        "academy": ROOT / "acceptance/screenshots/theme-v2/academy-dashboard.png",
        "command": ROOT / "acceptance/screenshots/theme-v2/command-dashboard.png",
    },
    "book": {
        "atlas": ROOT / "acceptance/screenshots/theme-v2/atlas-book.png",
        "academy": ROOT / "acceptance/screenshots/theme-v2/academy-book.png",
        "command": ROOT / "acceptance/screenshots/theme-v2/command-book.png",
    },
    "user": {
        "atlas": ROOT / "acceptance/screenshots/theme-v2/atlas-user.png",
        "academy": ROOT / "acceptance/screenshots/theme-v2/academy-user.png",
        "command": ROOT / "acceptance/screenshots/theme-v2/command-user.png",
    },
}
OLD_IMAGES = {
    "login": ROOT / "images/login.png",
    "dashboard": ROOT / "images/dashboard.png",
    "book": ROOT / "images/book.png",
    "user": ROOT / "images/reader.png",
}
AFTER_IMAGES = {
    "login": ROOT / "acceptance/screenshots/theme-v2/academy-login.png",
    "dashboard": ROOT / "acceptance/screenshots/theme-v2/academy-dashboard.png",
    "book": ROOT / "acceptance/screenshots/theme-v2/academy-book.png",
    "user": ROOT / "acceptance/screenshots/theme-v2/academy-user.png",
}


def color(value):
    return RGBColor.from_string(value)


def shape(slide, x, y, w, h, fill, rounded=False, line_color=None):
    item = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    item.fill.solid()
    item.fill.fore_color.rgb = color(fill)
    item.line.color.rgb = color(line_color or fill)
    if rounded:
        item.adjustments[0] = 0.08
    return item


def text(slide, value, x, y, w, h, size=16, fill=INK, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(fill)
    return box


def rule(slide, x1, y1, x2, y2, fill=LINE, width=1.2):
    item = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    item.line.color.rgb = color(fill)
    item.line.width = Pt(width)
    return item


def image_fit(slide, path, x, y, w, h, border=True):
    if not path.exists():
        raise FileNotFoundError(f"missing image: {path}")
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    nw, nh = iw * scale, ih * scale
    if border:
        shape(slide, x - 0.03, y - 0.03, w + 0.06, h + 0.06, WHITE, True, LINE)
    slide.shapes.add_picture(
        str(path), Inches(x + (w - nw) / 2), Inches(y + (h - nh) / 2),
        width=Inches(nw), height=Inches(nh),
    )


def card(slide, x, y, w, h, fill=WHITE, line_color=LINE):
    return shape(slide, x, y, w, h, fill, True, line_color)


def pill(slide, label, x, y, w, fill, text_fill=WHITE):
    shape(slide, x, y, w, 0.30, fill, True, fill)
    text(slide, label, x, y + 0.035, w, 0.20, 10, text_fill, True, PP_ALIGN.CENTER)


def bullet_list(slide, items, x, y, w, size=14, fill=INK, gap=0.43, dot=TEAL):
    for index, item in enumerate(items):
        yy = y + index * gap
        shape(slide, x, yy + 0.10, 0.08, 0.08, dot, True, dot)
        text(slide, item, x + 0.19, yy, w - 0.19, gap, size, fill)


def background(slide, dark=False):
    shape(slide, 0, 0, W, H, NAVY_2 if dark else CREAM)
    if dark:
        shape(slide, 0, 0, 0.18, H, PURPLE, False, PURPLE)


def heading(slide, kicker, title, subtitle="", dark=False):
    text(slide, kicker.upper(), 0.62, 0.34, 5.0, 0.22, 10, CYAN if dark else TEAL, True)
    text(slide, title, 0.62, 0.68, 12.0, 0.55, 26, WHITE if dark else NAVY, True)
    if subtitle:
        text(slide, subtitle, 0.64, 1.28, 12.0, 0.36, 12, "AABBD0" if dark else MUTED)


def footer(slide, number, dark=False):
    rule(slide, 0.62, 7.08, 12.70, 7.08, "314766" if dark else LINE, 0.8)
    text(slide, "图书馆管理系统 · 五项需求变化 · T4 汇报版", 0.64, 7.14, 6.5, 0.18, 8, "8EA3B8" if dark else MUTED)
    text(slide, f"{number:02d} / {TOTAL:02d}", 11.82, 7.14, 0.90, 0.18, 8, "8EA3B8" if dark else MUTED, align=PP_ALIGN.RIGHT)


def metric(slide, x, y, w, value, label, accent=TEAL, dark=False):
    fill = "14243E" if dark else WHITE
    card(slide, x, y, w, 1.10, fill, accent)
    shape(slide, x, y, 0.08, 1.10, accent, True, accent)
    text(slide, value, x + 0.22, y + 0.15, w - 0.30, 0.40, 25, WHITE if dark else NAVY, True)
    text(slide, label, x + 0.22, y + 0.68, w - 0.30, 0.22, 11, "AABBD0" if dark else MUTED)


def three_theme_slide(number, kicker, title, subtitle, page_key, conclusion):
    slide = prs.slides.add_slide(BLANK)
    background(slide)
    heading(slide, kicker, title, subtitle)
    labels = [
        ("atlas", "atlas · 完整左栏运营台", TEAL),
        ("academy", "academy · 顶部导航书院分栏", BRICK),
        ("command", "command · 窄轨道数字指挥舱", PURPLE),
    ]
    for index, (key, label, accent) in enumerate(labels):
        x = 0.55 + index * 4.10
        card(slide, x, 1.88, 3.84, 4.86, WHITE, accent)
        pill(slide, label, x + 0.18, 2.08, 3.46, accent)
        image_fit(slide, THEME_IMAGES[page_key][key], x + 0.16, 2.54, 3.52, 3.60)
        text(slide, {
            "atlas": "高密度左栏 + 顶部状态栏",
            "academy": "顶部横向导航 + 居中纸张内容",
            "command": "窄图标轨道 + 网格面板 + 悬浮工具栏",
        }[key], x + 0.18, 6.22, 3.46, 0.28, 11, INK, True, PP_ALIGN.CENTER)
    shape(slide, 0.76, 6.86, 11.80, 0.30, TEAL_LIGHT, True, TEAL_LIGHT)
    text(slide, "结论  ·  " + conclusion, 0.96, 6.91, 11.38, 0.20, 12, TEAL, True, PP_ALIGN.CENTER)
    footer(slide, number)


def before_after_tile(slide, x, y, title, old_path, new_path):
    card(slide, x, y, 5.92, 2.02, WHITE, LINE)
    text(slide, title, x + 0.18, y + 0.14, 2.0, 0.24, 14, NAVY, True)
    pill(slide, "原始", x + 3.70, y + 0.12, 0.62, "8796A1")
    pill(slide, "academy 最新", x + 4.38, y + 0.12, 1.15, BRICK)
    image_fit(slide, old_path, x + 0.18, y + 0.52, 2.68, 1.30, False)
    image_fit(slide, new_path, x + 3.05, y + 0.52, 2.68, 1.30, False)


def before_after_slide(number):
    slide = prs.slides.add_slide(BLANK)
    background(slide)
    heading(slide, "07 / BEFORE → LATEST", "保留原始主线：改造前 vs 最新 academy 默认风格", "原始仓库截图与最终真实 academy 截图并列；三主题差异在后续页面展开")
    before_after_tile(slide, 0.62, 1.82, "登录", OLD_IMAGES["login"], AFTER_IMAGES["login"])
    before_after_tile(slide, 6.80, 1.82, "Dashboard", OLD_IMAGES["dashboard"], AFTER_IMAGES["dashboard"])
    before_after_tile(slide, 0.62, 4.05, "图书", OLD_IMAGES["book"], AFTER_IMAGES["book"])
    before_after_tile(slide, 6.80, 4.05, "用户", OLD_IMAGES["user"], AFTER_IMAGES["user"])
    footer(slide, number)


# 01 cover
slide = prs.slides.add_slide(BLANK)
background(slide, True)
text(slide, "LIBRARY / DELIVERY REVIEW", 0.72, 0.62, 6.3, 0.25, 11, CYAN, True)
text(slide, "图书馆管理系统\n最终汇报", 0.72, 1.42, 7.2, 1.15, 34, WHITE, True)
text(slide, "五项需求 · 三套风格 · 一键验收", 0.75, 2.86, 7.0, 0.38, 21, "D4E1E2")
rule(slide, 0.76, 3.58, 5.28, 3.58, PURPLE, 2)
text(slide, "真实截图 · HTTP 黑盒 · 提交证据 · 边界说明", 0.76, 3.82, 7.2, 0.28, 14, "AABBD0")
card(slide, 8.72, 1.24, 3.58, 4.60, "14243E", PURPLE)
text(slide, "5", 9.18, 1.70, 2.65, 1.00, 76, CYAN, True, PP_ALIGN.CENTER)
text(slide, "项需求\n连续交付", 9.25, 2.98, 2.50, 0.78, 23, WHITE, True, PP_ALIGN.CENTER)
pill(slide, "T1 → T4", 9.62, 4.36, 1.82, PURPLE)
text(slide, "2026.07 · aef9188", 9.10, 5.18, 2.80, 0.24, 11, "AABBD0", align=PP_ALIGN.CENTER)
footer(slide, 1, True)

# 02 context and five requirements
slide = prs.slides.add_slide(BLANK)
background(slide)
heading(slide, "01 / CONTEXT", "五项需求：从“能用”走向“可验证”", "规则、数据、流程、权限与视觉共同形成一条可复查证据链")
card(slide, 0.62, 1.82, 4.05, 4.80, NAVY, NAVY)
text(slide, "原始断点", 0.96, 2.18, 2.0, 0.28, 20, WHITE, True)
bullet_list(slide, [
    "库存只有可借 / 不可借，无法表达馆藏规模",
    "借还操作分散，失败可能留下半成品状态",
    "逾期、续借和删除边界不统一",
    "页面信息层级不稳定，难以投影答辩",
], 0.98, 2.80, 3.25, 14, "D4E1E2", 0.72, CYAN)
card(slide, 4.98, 1.82, 7.74, 4.80, WHITE, LINE)
text(slide, "交付闭环", 5.34, 2.18, 2.0, 0.28, 20, NAVY, True)
items = [
    ("01", "库存数量", "totalCount / availableCount", TEAL),
    ("02", "逾期管理", "dueStatus / overdueDays", ORANGE),
    ("03", "流程后端化", "/circulation/*", NAVY),
    ("04", "用户管理", "新增 / 删除边界", RED),
    ("05", "三套视觉", "atlas / academy / command", PURPLE),
]
for index, (num, name, code, accent) in enumerate(items):
    x = 5.34 + (index % 3) * 2.35
    y = 2.88 + (index // 3) * 1.45
    card(slide, x, y, 2.08, 1.08, "F8FAFA", accent)
    text(slide, num, x + 0.14, y + 0.13, 0.44, 0.22, 14, accent, True)
    text(slide, name, x + 0.60, y + 0.12, 1.30, 0.24, 15, NAVY, True)
    text(slide, code, x + 0.14, y + 0.55, 1.78, 0.20, 10, MUTED)
text(slide, "验收口径：真实运行截图 + HTTP 黑盒 + 全量回归，不把“代码已写”当作“能力已完成”。", 5.34, 5.96, 6.55, 0.40, 15, TEAL, True)
footer(slide, 2)

# 03 inventory
slide = prs.slides.add_slide(BLANK)
background(slide)
heading(slide, "02 / REQUIREMENT 1", "库存从“状态”升级为“数量链路”", "页面展示与数据库约束共用 totalCount、availableCount 和已借出数")
card(slide, 0.62, 1.88, 7.20, 4.70, WHITE, LINE)
text(slide, "一次借书如何改变数据", 0.98, 2.22, 3.5, 0.28, 19, NAVY, True)
stages = [("借前", "总数 1", "可借 1", TEAL_LIGHT, TEAL), ("借书成功", "库存 -1", "当前借阅 +1", BRICK_LIGHT, ORANGE), ("借后", "总数 1", "可借 0", "FCE9E7", RED)]
for index, (label, a, b, fill, accent) in enumerate(stages):
    x = 1.0 + index * 2.13
    card(slide, x, 2.94, 1.76, 1.42, fill, LINE)
    text(slide, label, x, 3.15, 1.76, 0.22, 13, accent, True, PP_ALIGN.CENTER)
    text(slide, a, x, 3.53, 1.76, 0.22, 14, INK, True, PP_ALIGN.CENTER)
    text(slide, b, x, 3.88, 1.76, 0.20, 11, MUTED, align=PP_ALIGN.CENTER)
    if index < 2:
        rule(slide, x + 1.80, 3.64, x + 2.03, 3.64, accent, 1.6)
text(slide, "后端约束", 1.0, 4.86, 1.3, 0.22, 12, MUTED, True)
bullet_list(slide, [
    "新增：totalCount 必须为正整数，availableCount 由后端初始化",
    "编辑：总数不得小于当前已借出数，并发修改串行化",
    "重复借阅：一成一败，库存与当前借阅不穿透",
], 1.0, 5.18, 6.25, 13, INK, 0.38)
card(slide, 8.10, 1.88, 4.62, 4.70, NAVY, NAVY)
text(slide, "真实黑盒", 8.48, 2.22, 2.0, 0.25, 13, CYAN, True)
text(slide, "9/9", 8.48, 2.68, 2.2, 0.70, 44, WHITE, True)
text(slide, "库存初始化 · 借还变化\n非法总数 · 并发编辑\n刷新后查询仍一致", 8.52, 3.54, 3.25, 0.85, 15, "D4E1E2", True)
pill(slide, "verify_inventory_http.py", 8.48, 5.27, 2.62, TEAL)
footer(slide, 3)

# 04 overdue
slide = prs.slides.add_slide(BLANK)
background(slide)
heading(slide, "03 / REQUIREMENT 2", "逾期状态成为可执行的借阅规则", "后端按 Asia/Shanghai 计算自然日边界，状态、天数与限制同时返回")
card(slide, 0.62, 1.88, 7.20, 4.70, WHITE, LINE)
text(slide, "状态机", 0.98, 2.22, 1.2, 0.28, 19, NAVY, True)
states = [("NORMAL", "> 3 天", TEAL, TEAL_LIGHT), ("DUE_SOON", "0–3 天", ORANGE, BRICK_LIGHT), ("OVERDUE", "早于今天", RED, "FCE9E7")]
for index, (label, desc, accent, fill) in enumerate(states):
    x = 1.02 + index * 2.13
    card(slide, x, 2.96, 1.78, 1.20, fill, LINE)
    text(slide, label, x, 3.18, 1.78, 0.22, 14, accent, True, PP_ALIGN.CENTER)
    text(slide, desc, x, 3.58, 1.78, 0.20, 11, INK, align=PP_ALIGN.CENTER)
    if index < 2:
        rule(slide, x + 1.82, 3.55, x + 2.03, 3.55, accent, 1.8)
text(slide, "逾期后", 1.02, 4.76, 1.1, 0.22, 12, MUTED, True)
text(slide, "借新书 ✕   续借逾期图书 ✕   归还后自动解除限制 ✓", 2.10, 4.72, 5.35, 0.28, 15, INK, True)
text(slide, "规则可被筛选、纠正、复验，而不是只在页面上显示颜色。", 1.02, 5.47, 6.20, 0.34, 15, TEAL, True)
card(slide, 8.10, 1.88, 4.62, 4.70, NAVY, NAVY)
text(slide, "真实黑盒", 8.48, 2.22, 2.0, 0.25, 13, CYAN, True)
text(slide, "6/6", 8.48, 2.68, 2.2, 0.70, 44, WHITE, True)
text(slide, "+30 天生成 · 昨天=1天逾期\n借阅/续借限制 · overdueOnly\n归还后自动解除", 8.52, 3.54, 3.35, 0.88, 15, "D4E1E2", True)
pill(slide, "verify_overdue_http.py", 8.48, 5.27, 2.40, RED)
footer(slide, 4)

# 05 circulation
slide = prs.slides.add_slide(BLANK)
background(slide)
heading(slide, "04 / REQUIREMENT 3", "借书、还书、续借收敛为事务业务接口", "前端只提交 readerId + isbn；后端一次完成校验、库存、日期与三表更新")
card(slide, 0.62, 1.88, 12.10, 3.86, WHITE, LINE)
nodes = [("页面点击", "readerId + isbn", 1.0, BRICK_LIGHT, ORANGE), ("Controller", "/circulation/*", 3.72, TEAL_LIGHT, TEAL), ("Service", "@Transactional", 6.44, "E8ECF7", PURPLE), ("三表一致", "book · current · history", 9.16, TEAL_LIGHT, TEAL)]
for index, (a, b, x, fill, accent) in enumerate(nodes):
    card(slide, x, 2.78, 2.20, 1.50, fill, LINE)
    text(slide, a, x + 0.1, 3.08, 2.0, 0.22, 15, accent, True, PP_ALIGN.CENTER)
    text(slide, b, x + 0.1, 3.50, 2.0, 0.22, 11, INK, align=PP_ALIGN.CENTER)
    if index < 3:
        rule(slide, x + 2.22, 3.52, x + 2.62, 3.52, accent, 2)
text(slide, "失败路径数据不变：库存不足、重复借阅、重复还书、第二次续借都会被后端拒绝。", 1.0, 4.86, 10.9, 0.30, 15, TEAL, True, PP_ALIGN.CENTER)
metric(slide, 0.72, 6.00, 2.45, "3", "业务接口", TEAL)
metric(slide, 3.37, 6.00, 2.45, "1", "页面一次写请求", ORANGE)
metric(slide, 6.02, 6.00, 2.45, "6/6", "借还续黑盒", PURPLE)
metric(slide, 8.67, 6.00, 3.50, "PASS", "失败路径数据不变", RED)
footer(slide, 5)

# 06 user
slide = prs.slides.add_slide(BLANK)
background(slide)
heading(slide, "05 / REQUIREMENT 4", "用户管理补上“新增—约束—删除”闭环", "管理员可新增普通读者；有未归还图书时禁止删除，归还后才可清理")
card(slide, 0.62, 1.88, 7.35, 4.70, WHITE, LINE)
text(slide, "管理员操作链路", 0.98, 2.22, 2.8, 0.28, 19, NAVY, True)
chain = [("新增", "role=2", TEAL, TEAL_LIGHT), ("重复名", "拒绝", ORANGE, BRICK_LIGHT), ("有借阅", "禁止删", RED, "FCE9E7"), ("归还后", "成功删", TEAL, TEAL_LIGHT)]
for index, (a, b, accent, fill) in enumerate(chain):
    x = 0.98 + index * 1.70
    card(slide, x, 3.04, 1.42, 1.38, fill, LINE)
    text(slide, a, x, 3.30, 1.42, 0.22, 13, accent, True, PP_ALIGN.CENTER)
    text(slide, b, x, 3.70, 1.42, 0.20, 11, INK, align=PP_ALIGN.CENTER)
    if index < 3:
        rule(slide, x + 1.44, 3.70, x + 1.60, 3.70, accent, 1.5)
bullet_list(slide, [
    "operatorId 对应管理员才可新增 / 删除",
    "管理员目标与旧批量删除路径明确拒绝",
    "正式 HTTP 黑盒覆盖借阅中删除保护",
], 0.98, 4.92, 6.3, 13, INK, 0.38)
card(slide, 8.22, 1.88, 4.50, 4.70, NAVY, NAVY)
text(slide, "真实黑盒", 8.56, 2.22, 2.0, 0.25, 13, CYAN, True)
text(slide, "7/7", 8.56, 2.68, 2.2, 0.70, 44, WHITE, True)
text(slide, "新增可登录 · 重复名拒绝\n有借阅禁止删 · 归还后可删", 8.60, 3.54, 3.35, 0.62, 15, "D4E1E2", True)
pill(slide, "uk_user_username", 8.56, 4.55, 2.25, TEAL)
pill(slide, "verify_user_management_http.py", 8.56, 5.02, 3.12, PURPLE)
footer(slide, 6)

# 07 before/after baseline
before_after_slide(7)

# 08 visual overview, using final real screenshots
slide = prs.slides.add_slide(BLANK)
background(slide)
heading(slide, "08 / REQUIREMENT 5", "最新三风格总览：同一业务，不同工作方式", "默认 academy；登录前和登录后都可切换，localStorage 持久化，切换无需刷新")
overview = [("atlas", "完整左栏运营台", TEAL, "完整侧栏 · 顶部状态栏 · 高信息密度"), ("academy", "顶部导航书院分栏", BRICK, "顶部横向导航 · 纸张留白 · 居中内容"), ("command", "窄轨道数字指挥舱", PURPLE, "窄图标轨道 · 网格面板 · 悬浮工具栏")]
for index, (key, label, accent, desc) in enumerate(overview):
    x = 0.55 + index * 4.10
    card(slide, x, 1.88, 3.84, 4.92, WHITE, accent)
    pill(slide, label, x + 0.18, 2.08, 3.46, accent)
    image_fit(slide, THEME_IMAGES["login"][key], x + 0.16, 2.54, 3.52, 2.20)
    text(slide, desc, x + 0.25, 5.06, 3.34, 0.54, 13, INK, True, PP_ALIGN.CENTER)
    text(slide, {"atlas": "高密度运营", "academy": "默认书院阅读风", "command": "实时指挥舱"}[key], x + 0.25, 5.86, 3.34, 0.28, 12, accent, True, PP_ALIGN.CENTER)
text(slide, "三套框架共享业务页面与 API，不复制三份路由或保存逻辑。", 0.80, 6.86, 11.70, 0.22, 13, TEAL, True, PP_ALIGN.CENTER)
footer(slide, 8)

# 09–12 exact final screenshots, three columns each
three_theme_slide(9, "09 / FINAL SCREENSHOTS", "Login：登录前即可选择三种工作方式", "同一账号入口，三种视觉框架；登录后 Header 保持同一选择器", "login", "academy 是首次访问默认值，切换只更新主题状态，不触发整页刷新。")
three_theme_slide(10, "10 / FINAL SCREENSHOTS", "Dashboard：指标与图表随框架重新编排", "Atlas 高密度卡片、Academy 杂志式留白、Command 网格控制台", "dashboard", "差异来自布局与信息层级，而不只是颜色。")
three_theme_slide(11, "11 / FINAL SCREENSHOTS", "Book：筛选、操作区与数据表按主题重排", "三主题均使用同一业务 methods；库存、借阅、编辑行为保持不变", "book", "宽表只在业务表格容器内滚动，页面本身无无意义横向溢出。")
three_theme_slide(12, "12 / FINAL SCREENSHOTS", "User：干净中文数据证明真实页面验收", "截图使用临时中文读者按姓名筛选后的真实 User 页面，非静态伪造", "user", "三套布局都只显示一条中文读者记录，历史乱码与 verify_reader 数据不进入交付证据。")

# 13 one-click backend demo
slide = prs.slides.add_slide(BLANK)
background(slide, True)
heading(slide, "13 / ONE-CLICK BACKEND DEMO", "一条命令跑完四组后端黑盒", "脚本位置：/home/nianjiu/下载/617/library/demo_backend_verification.sh", True)
card(slide, 0.72, 1.86, 5.30, 4.74, "14243E", PURPLE)
text(slide, "现场命令", 1.08, 2.22, 2.0, 0.25, 14, CYAN, True)
shape(slide, 1.08, 2.72, 4.56, 0.84, NAVY_2, True, "314766")
text(slide, "./demo_backend_verification.sh", 1.28, 2.98, 4.15, 0.25, 15, WHITE, True)
text(slide, "先检查 /dashboard；单项失败继续；末尾汇总真实退出码。", 1.08, 4.02, 4.30, 0.55, 15, "D4E1E2")
pill(slide, "HTTP API 临时数据", 1.08, 5.18, 1.85, TEAL)
text(slide, "各 Python 黑盒自行创建并清理临时读者、图书、借阅与历史；编排脚本不直连数据库。", 1.08, 5.72, 4.42, 0.48, 11, "AABBD0")
card(slide, 6.42, 1.86, 6.26, 4.74, WHITE, LINE)
text(slide, "TOTAL PASS", 6.82, 2.24, 3.0, 0.36, 26, NAVY, True)
rows = [("circulation", "6/6", TEAL), ("inventory", "9/9", ORANGE), ("overdue", "6/6", RED), ("user management", "7/7", PURPLE)]
for index, (label, value, accent) in enumerate(rows):
    y = 2.98 + index * 0.72
    shape(slide, 6.84, y, 5.36, 0.52, "F8FAFA", True, LINE)
    text(slide, label, 7.06, y + 0.13, 3.15, 0.20, 13, INK, True)
    text(slide, value, 10.82, y + 0.10, 1.05, 0.25, 18, accent, True, PP_ALIGN.RIGHT)
text(slide, "默认运行、显式 ADMIN_PASSWORD=123456、失败继续分支均已实测。", 6.84, 6.10, 5.25, 0.30, 12, TEAL, True)
footer(slide, 13, True)

# 14 evidence and commit distinction
slide = prs.slides.add_slide(BLANK)
background(slide)
heading(slide, "14 / EVIDENCE", "提交与验收证据：原五项需求 + 第二版视觉增强", "aef9188 是最新三风格提交；PPT 本身作为独立 T4 汇报交付物，不再制造提交")
card(slide, 0.62, 1.86, 7.18, 4.76, WHITE, LINE)
text(slide, "原五项需求主线", 0.98, 2.22, 3.0, 0.28, 19, NAVY, True)
commits = [("1745b11", "后端化借书 / 还书 / 续借"), ("eca083e", "库存数量与数据库约束"), ("23401f3", "逾期管理与借阅限制"), ("322e1ce", "管理员读者新增 / 删除"), ("48c645d", "前端视觉与交互重构")]
for index, (commit, label) in enumerate(commits):
    y = 2.92 + index * 0.55
    pill(slide, f"T{index + 1}", 1.0, y, 0.44, TEAL)
    text(slide, commit, 1.58, y + 0.035, 1.10, 0.20, 11, NAVY, True)
    text(slide, label, 2.88, y + 0.035, 4.35, 0.20, 13, INK)
card(slide, 8.18, 1.86, 4.54, 4.76, NAVY, NAVY)
text(slide, "第二版视觉增强", 8.56, 2.22, 3.0, 0.28, 19, CYAN, True)
text(slide, "aef9188", 8.56, 2.78, 2.5, 0.42, 26, WHITE, True)
text(slide, "增加 atlas / academy / command\n主题状态、布局分支、响应式\n真实 Chromium 截图与 browser-check.json", 8.56, 3.54, 3.45, 1.00, 16, "D4E1E2", True)
pill(slide, "14 / 14 页面证据", 8.56, 5.25, 2.08, PURPLE)
text(slide, "验收：verify、verify-full、四项黑盒、Playwright、build、Docker", 8.56, 5.82, 3.42, 0.42, 11, "AABBD0")
footer(slide, 14)

# 15 boundaries
slide = prs.slides.add_slide(BLANK)
background(slide)
heading(slide, "15 / BOUNDARIES", "验收边界：通过的写清楚，未完成的不包装", "答辩时明确哪些是本次证据，哪些仍需负责人后续处理")
card(slide, 0.62, 1.88, 5.92, 4.74, NAVY, NAVY)
text(slide, "已验证", 1.0, 2.25, 1.5, 0.28, 19, CYAN, True)
bullet_list(slide, [
    "四组 HTTP 黑盒：6/6、9/9、6/6、7/7",
    "三主题真实 Login / Dashboard / Book / User 截图",
    "主题默认 academy、localStorage、无刷新切换",
    "前端 build、Docker frontend、全量 verify",
], 1.0, 2.90, 4.80, 15, "D4E1E2", 0.68, CYAN)
card(slide, 6.82, 1.88, 5.92, 4.74, WHITE, LINE)
text(slide, "明确保留的边界", 7.20, 2.25, 2.8, 0.28, 19, NAVY, True)
bullet_list(slide, [
    "普通读者可查看全部读者列表的既有权限 WARN",
    "未新增统一 token 鉴权拦截器",
    "PPT 由负责人审查，不在本次任务 commit/push",
    "不把未执行的人工写操作包装成已完成",
], 7.20, 2.90, 4.80, 15, INK, 0.68, BRICK)
footer(slide, 15)

# 16 closing
slide = prs.slides.add_slide(BLANK)
background(slide, True)
heading(slide, "16 / CLOSING", "结论：五项需求形成一条可追踪证据链", "从后端规则、数据一致性到三套视觉框架，最终交付可运行、可演示、可复查", True)
card(slide, 0.72, 1.92, 5.80, 4.66, "14243E", PURPLE)
text(slide, "答辩结论", 1.10, 2.30, 2.0, 0.28, 20, CYAN, True)
text(slide, "规则落后端\n数据可核对\n页面可切换\n验收可复现", 1.10, 3.00, 3.0, 1.65, 25, WHITE, True)
pill(slide, "TOTAL PASS", 1.10, 5.48, 1.72, PURPLE)
card(slide, 6.86, 1.92, 5.76, 4.66, "14243E", PURPLE)
text(slide, "下一步", 7.24, 2.30, 1.5, 0.28, 20, "F4B187", True)
bullet_list(slide, [
    "负责人审查最终 PPT 与 12 张真实截图",
    "如需现场演示，执行 ./demo_backend_verification.sh",
    "保留既有权限 WARN 边界，不扩大本次范围",
], 7.24, 3.02, 4.62, 17, "D4E1E2", 0.75, ORANGE)
text(slide, "谢谢", 7.24, 5.50, 4.4, 0.46, 28, WHITE, True)
footer(slide, 16, True)


required_images = [*OLD_IMAGES.values(), *AFTER_IMAGES.values()]
required_images.extend(path for page in THEME_IMAGES.values() for path in page.values())
missing = [str(path) for path in required_images if not path.exists()]
if missing:
    raise FileNotFoundError("missing required images:\n" + "\n".join(missing))

prs.save(OUT)
print(f"saved {OUT} ({len(prs.slides)} slides, {W}x{H})")
